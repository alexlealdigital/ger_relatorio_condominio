"""
Gerador de Apresentações Premium v2.0 — Gestor Financeiro de Condomínios
Design moderno 16:9, tema navy, gráficos estilizados.

Narrativa (construída para fundamentar a proposta de reajuste):
  1. Capa
  2. Panorama do ano atual (previsto × realizado — o argumento)
  3. Evolução mensal das despesas contra a previsão aprovada
  4. Composição das despesas por categoria (donut)
  5. Maiores despesas (ranking)
  6. Inadimplência (opcional — pressão sobre o caixa)
  7. Previsão do próximo ano (composição do rateio)
  8. A proposta (taxa atual → taxa ideal, reajuste em destaque)
  9. Encerramento (proposta para aprovação)
  + Apêndice: detalhamento das rubricas

Compatível com os dois formatos do analisador (balanual e legado).
Interface: PowerPointGenerator(analyzer, output_path, extras=None).generate()
  extras = {'inadimplencia_total': float, 'acordos_receber': float}
`output_path` aceita caminho ou stream (BytesIO).
"""

import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib import font_manager

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Paleta (identidade navy premium) ────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x1F, 0x44)
NAVY_2    = RGBColor(0x11, 0x22, 0x55)
BLUE      = RGBColor(0x1A, 0x56, 0xDB)
BLUE_LT   = RGBColor(0x6E, 0x9B, 0xFF)
ICE       = RGBColor(0xEB, 0xF0, 0xFB)
PEACH     = RGBColor(0xFB, 0xE3, 0xCC)
ICE_2     = RGBColor(0xF5, 0xF7, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x0E, 0x9F, 0x6E)
AMBER     = RGBColor(0xF5, 0xA6, 0x23)
RED       = RGBColor(0xE8, 0x50, 0x3A)
TEXT      = RGBColor(0x1B, 0x26, 0x3B)
MUTED     = RGBColor(0x5A, 0x6B, 0x8C)
MUTED_ICE = RGBColor(0xAF, 0xC2, 0xE8)

# Paleta hex para matplotlib
MPL_SEQ = ['#1A56DB', '#6E9BFF', '#0E9F6E', '#F5A623', '#8A63D2', '#E8503A',
           '#2FA8C9', '#B0B9CC']
FONT = 'Cambria'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _fmt_brl(v, dec=2):
    if v is None:
        return '—'
    s = f'{v:,.{dec}f}'.replace(',', 'X').replace('.', ',').replace('X', '.')
    return f'R$ {s}'


def _fmt_pct(v, dec=1, sinal=False):
    if v is None:
        return '—'
    s = f'{v:+.{dec}f}' if sinal else f'{v:.{dec}f}'
    return s.replace('.', ',') + '%'


class PowerPointGenerator:
    def __init__(self, analyzer, output_path, extras=None):
        self.r = analyzer.analysis_results
        self.b = self.r.get('balanual') or {}
        self.output_path = output_path
        self.extras = extras or {}
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════════
    # PRIMITIVAS
    # ══════════════════════════════════════════════════════════════════
    def _slide(self, bg=WHITE):
        s = self.prs.slides.add_slide(self._blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
        return s

    def _text(self, slide, x, y, w, h, texto, size=16, color=TEXT, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        linhas = texto.upper().split('\n')
        for i, linha in enumerate(linhas):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if spacing:
                p.space_after = spacing
            run = p.add_run()
            run.text = linha
            f = run.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = bold
            f.color.rgb = color
        return tb

    def _rich(self, slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
        """runs: lista de (texto, size, color, bold) em um único parágrafo."""
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = align
        for texto, size, color, bold in runs:
            run = p.add_run()
            run.text = texto.upper()
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return tb

    def _card(self, slide, x, y, w, h, fill=ICE_2, line=None, radius=0.09):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
            shp.line.width = Pt(1)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _kpi(self, slide, x, y, w, h, rotulo, valor, sub=None,
             cor_valor=TEXT, fill=ICE_2, size_valor=30):
        self._card(slide, x, y, w, h, fill=fill)
        pad = Inches(0.22)
        self._text(slide, x + pad, y + Inches(0.18), w - 2 * pad, Inches(0.3),
                   rotulo.upper(), size=11, color=MUTED, bold=True)
        self._text(slide, x + pad, y + Inches(0.52), w - 2 * pad, Inches(0.62),
                   valor, size=size_valor, color=cor_valor, bold=True)
        if sub:
            self._text(slide, x + pad, y + h - Inches(0.42), w - 2 * pad,
                       Inches(0.32), sub, size=10.5, color=MUTED)

    def _titulo(self, slide, texto, sub=None, dark=False):
        self._logo_canto(slide)
        cor = WHITE if dark else NAVY
        # título deixa folga à direita para a logo
        self._text(slide, Inches(0.6), Inches(0.42), Inches(10.4), Inches(0.7),
                   texto, size=30, color=cor, bold=True)
        if sub:
            self._text(slide, Inches(0.6), Inches(1.06), Inches(10.4),
                       Inches(0.36), sub, size=14,
                       color=MUTED_ICE if dark else MUTED)

    def _logo(self, slide, x, y, altura, direita=False):
        """Insere a logo (se fornecida) preservando a proporção.
        Com direita=True, alinha ao canto direito (x vira a margem)."""
        path = self.extras.get('logo_path')
        if not path:
            return
        try:
            from PIL import Image as _Img
            with _Img.open(path) as im:
                prop = im.width / im.height
            larg = Emu(int(altura * prop))
            px = (SLIDE_W - larg - x) if direita else x
            slide.shapes.add_picture(path, px, y, height=altura, width=larg)
        except Exception:
            pass

    def _logo_canto(self, slide):
        """Logo no canto superior DIREITO dos slides de conteúdo (padrão
        do cliente)."""
        path = self.extras.get('logo_path')
        if not path:
            return 0
        try:
            from PIL import Image as _Img
            with _Img.open(path) as im:
                prop = im.width / im.height
            alt = Inches(0.8)
            larg = Emu(int(alt * prop))
            slide.shapes.add_picture(path, SLIDE_W - larg - Inches(0.6),
                                     Inches(0.38), height=alt, width=larg)
            return 0
        except Exception:
            return 0

    def _rodape(self, slide, dark=False):
        nome = self.r.get('nome_condominio') or 'PREVISÃO ORÇAMENTÁRIA'
        self._text(slide, Inches(0.6), Inches(7.08), Inches(12.1), Inches(0.3),
                   nome, size=9, color=MUTED_ICE if dark else RGBColor(0xB6, 0xC0, 0xD4))

    def _add_chart(self, slide, fig, x, y, w, h=None):
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=200, transparent=True,
                    bbox_inches='tight')
        plt.close(fig)
        buf.seek(0)
        if h:
            slide.shapes.add_picture(buf, x, y, width=w, height=h)
        else:
            slide.shapes.add_picture(buf, x, y, width=w)

    @staticmethod
    def _mpl_base(ax):
        for sp in ('top', 'right', 'left'):
            ax.spines[sp].set_visible(False)
        ax.spines['bottom'].set_color('#D5DCEA')
        ax.tick_params(colors='#5A6B8C', labelsize=11, length=0)
        ax.yaxis.grid(True, color='#E4E9F4', linewidth=0.9)
        ax.set_axisbelow(True)

    # ══════════════════════════════════════════════════════════════════
    # SLIDES
    # ══════════════════════════════════════════════════════════════════
    def _slide_capa(self):
        s = self._slide(bg=WHITE)
        # círculo pêssego à direita (cores da logo) — padrão do cliente
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(-1.9),
                                  Inches(6.2), Inches(6.2))
        circ.fill.solid()
        circ.fill.fore_color.rgb = PEACH
        circ.line.fill.background()
        circ.shadow.inherit = False
        nome = self.r.get('nome_condominio') or 'CONDOMÍNIO'
        ano = self.r.get('ano_proximo') or ''
        data = self.r.get('data_assembleia')

        self._text(s, Inches(0.9), Inches(0.85), Inches(8.0), Inches(0.45),
                   'ASSEMBLEIA GERAL ORDINÁRIA', size=16, color=NAVY,
                   bold=True)
        if data:
            self._text(s, Inches(0.9), Inches(1.35), Inches(8.0), Inches(0.38),
                       f'DATA: {data}', size=12.5, color=TEXT, bold=True)
        self._logo(s, Inches(10.35), Inches(0.7), altura=Inches(1.05))
        self._text(s, Inches(0.9), Inches(2.85), Inches(11.0), Inches(2.0),
                   f'Previsão\nOrçamentária {ano}', size=52, color=NAVY,
                   bold=True)
        self._text(s, Inches(0.9), Inches(5.35), Inches(11.0), Inches(0.5),
                   nome, size=20, color=TEXT, bold=True)

    def _slides_edital(self):
        """Edital de convocação após a capa: texto legível (preferido) ou imagem."""
        texto = (self.extras.get('edital_texto') or '').strip()
        if texto:
            # divide em blocos de ~2300 caracteres respeitando parágrafos
            paragrafos = [p.strip() for p in texto.split('\n') if p.strip()]
            blocos, atual = [], ''
            for p in paragrafos:
                if atual and len(atual) + len(p) > 2300:
                    blocos.append(atual)
                    atual = p
                else:
                    atual = f'{atual}\n{p}' if atual else p
            if atual:
                blocos.append(atual)
            for i, bloco in enumerate(blocos[:3]):
                s = self._slide(bg=WHITE)
                sufixo = f' ({i+1}/{len(blocos)})' if len(blocos) > 1 else ''
                self._titulo(s, f'Edital de Convocação{sufixo}')
                self._text(s, Inches(0.9), Inches(1.7), Inches(11.5),
                           Inches(5.2), bloco, size=15, color=TEXT,
                           spacing=Pt(8))
                self._rodape(s)
            return
        paths = self.extras.get('edital_paths') or []
        for path in paths[:3]:
            s = self._slide(bg=WHITE)
            self._titulo(s, 'Edital de Convocação')
            try:
                from PIL import Image as _Img
                with _Img.open(path) as im:
                    prop = im.width / im.height
                alt_max = Inches(5.6)
                larg_max = Inches(11.0)
                if larg_max / prop <= alt_max:
                    larg, alt = larg_max, Emu(int(larg_max / prop))
                else:
                    alt, larg = alt_max, Emu(int(alt_max * prop))
                s.shapes.add_picture(path, Emu(int((SLIDE_W - larg) / 2)),
                                     Inches(1.55), width=larg, height=alt)
            except Exception:
                pass
            self._rodape(s)

    def _slide_panorama(self):
        b = self.b
        if not b:
            return
        s = self._slide()
        ano = self.r.get('ano_atual') or ''
        sub_pan = ('Execução orçamentária do exercício — previsto × realizado'
                   if b.get('prev_aprovada_mensal') is not None
                   else 'Execução do exercício e projeção para o próximo ano')
        self._titulo(s, f'Panorama {ano}', sub_pan)

        y, h = Inches(1.85), Inches(1.75)
        w, gap = Inches(3.87), Inches(0.28)
        x0 = Inches(0.6)

        desvio = b.get('desvio_pct')
        tem_aprovada = b.get('prev_aprovada_mensal') is not None
        if tem_aprovada:
            cor_desvio = RED if (desvio or 0) > 0 else GREEN
            self._kpi(s, x0, y, w, h, 'Previsão aprovada / mês',
                      _fmt_brl(b.get('prev_aprovada_mensal'), 2),
                      sub=f'Orçamento aprovado para {ano}', size_valor=26)
            self._kpi(s, x0 + (w + gap), y, w, h, 'Realizado médio / mês',
                      _fmt_brl(b.get('realizado_medio_mensal'), 2),
                      sub=f'Total no ano: {_fmt_brl(b.get("realizado_total_anual"), 2)}',
                      size_valor=26)
            self._kpi(s, x0 + 2 * (w + gap), y, w, h, 'Desvio sobre o aprovado',
                      _fmt_pct(desvio, 1, sinal=True),
                      sub='Realizado acima da previsão' if (desvio or 0) > 0
                          else 'Realizado dentro da previsão',
                      cor_valor=cor_desvio, fill=ICE)
        else:
            ano_prox = self.r.get('ano_proximo') or ''
            reaj = self.r.get('reajuste_pct')
            og = b.get('taxa_atual_origem')
            sub_reaj = ('Sobre o total provisionado no período anterior'
                        if og in ('planilha', 'aprovada')
                        else 'Sobre a taxa vigente informada')
            self._kpi(s, x0, y, w, h, 'Realizado médio / mês',
                      _fmt_brl(b.get('realizado_medio_mensal'), 2),
                      sub=f'Total no ano: {_fmt_brl(b.get("realizado_total_anual"), 2)}',
                      size_valor=26)
            self._kpi(s, x0 + (w + gap), y, w, h, f'Previsão {ano_prox} / mês',
                      _fmt_brl(b.get('prev_next_mensal'), 2),
                      sub='Despesas projetadas para o próximo exercício',
                      size_valor=26)
            self._kpi(s, x0 + 2 * (w + gap), y, w, h, 'Reajuste projetado',
                      _fmt_pct(reaj, 2, sinal=True) if reaj is not None else '—',
                      sub=sub_reaj,
                      cor_valor=RED if (reaj or 0) > 0 else GREEN, fill=ICE)

        # mensagem-síntese
        if tem_aprovada and desvio and desvio > 0:
            msg = (f'As despesas reais superaram a previsão aprovada em '
                   f'{_fmt_pct(desvio)} na média mensal. Manter a taxa atual '
                   f'significa operar abaixo do custo real do condomínio.')
        elif tem_aprovada:
            msg = ('As despesas reais se mantiveram dentro da previsão '
                   'aprovada no exercício.')
        else:
            msg = (f'A despesa média mensal apurada em {ano} é a base da '
                   f'previsão orçamentária do próximo exercício e do valor '
                   f'da nova taxa condominial.')
        self._card(s, Inches(0.6), Inches(4.05), Inches(12.13), Inches(1.5),
                   fill=NAVY)
        self._text(s, Inches(1.0), Inches(4.32), Inches(11.3), Inches(1.0),
                   msg, size=17, color=WHITE, bold=True)

        uni = self.r.get('num_unidades')
        itens = b.get('num_itens')
        partes = ([f'{uni} unidades'] if uni else []) +                  ([f'{itens} rubricas de despesa analisadas'] if itens else [])
        if partes:
            self._text(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.4),
                       '  ·  '.join(partes), size=12.5, color=MUTED)
        self._rodape(s)

    def _slide_evolucao(self):
        b = self.b
        serie = b.get('serie_mensal') or []
        if not serie:
            return
        s = self._slide()
        ano = self.r.get('ano_atual') or ''
        prev_hdr = b.get('prev_aprovada_mensal')
        sub_ev = ('Barras: realizado no mês · Linha tracejada: previsão mensal aprovada'
                  if prev_hdr else 'Despesa realizada em cada mês do exercício')
        self._titulo(s, f'Evolução mensal das despesas — {ano}', sub_ev)

        prev = b.get('prev_aprovada_mensal')
        rotulos = [p['rotulo'].upper() for p in serie]
        valores = [p['valor'] / 1000.0 for p in serie]
        acima = sum(1 for p in serie if prev and p['valor'] > prev)

        fig, ax = plt.subplots(figsize=(11.6, 4.1))
        cores = ['#E8503A' if (prev and v * 1000 > prev) else '#1A56DB'
                 for v in valores]
        barras_ev = ax.bar(rotulos, valores, color=cores, width=0.62, zorder=3)
        vmax = max(valores) if valores else 1
        for barra, p in zip(barras_ev, serie):
            ax.text(barra.get_x() + barra.get_width() / 2,
                    barra.get_height() + vmax * 0.022,
                    f"{p['valor'] / 1000:,.1f}".replace(',', 'X')
                        .replace('.', ',').replace('X', '.'),
                    ha='center', va='bottom', fontsize=9.5,
                    color='#1B263B', fontweight='bold')
        ax.set_ylim(0, vmax * 1.14)
        if prev:
            ax.axhline(prev / 1000.0, color='#0A1F44', linewidth=1.6,
                       linestyle=(0, (6, 4)), zorder=4)
        self._mpl_base(ax)
        ax.set_ylabel('R$ mil / mês', fontsize=11, color='#5A6B8C')
        fig.tight_layout()
        self._add_chart(s, fig, Inches(0.55), Inches(1.6), Inches(12.2))

        if prev:
            self._rich(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.55), [
                ('Em ', 14, TEXT, False),
                (f'{acima} de {len(serie)} meses ', 14, RED, True),
                ('as despesas superaram a previsão aprovada.', 14, TEXT, False),
            ])
        self._rodape(s)

    def _slide_categorias(self):
        b = self.b
        cats = b.get('categorias') or []
        if not cats:
            return
        s = self._slide()
        self._titulo(s, 'Para onde vai o dinheiro',
                     'Composição da despesa média mensal por categoria')

        total = sum(c['mensal'] for c in cats)
        fig, ax = plt.subplots(figsize=(5.1, 5.1))
        vals = [c['mensal'] for c in cats]
        cores = MPL_SEQ[:len(cats)]
        ax.pie(vals, colors=cores, startangle=90, counterclock=False,
               wedgeprops={'width': 0.34, 'edgecolor': 'white', 'linewidth': 2})
        ax.text(0, 0.08, _fmt_brl(total, 0), ha='center', va='center',
                fontsize=17, fontweight='bold', color='#0A1F44')
        ax.text(0, -0.14, 'média / mês', ha='center', va='center',
                fontsize=10, color='#5A6B8C')
        self._add_chart(s, fig, Inches(0.7), Inches(1.7), Inches(4.9))

        # legenda em cards à direita
        y = Inches(1.75)
        for i, c in enumerate(cats[:7]):
            cor = RGBColor(*(int(MPL_SEQ[i][j:j+2], 16) for j in (1, 3, 5)))
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.15), y + Inches(0.06),
                                     Inches(0.18), Inches(0.18))
            dot.fill.solid()
            dot.fill.fore_color.rgb = cor
            dot.line.fill.background()
            dot.shadow.inherit = False
            self._text(s, Inches(6.5), y, Inches(4.4), Inches(0.32),
                       c['nome'], size=14, color=TEXT, bold=True)
            self._text(s, Inches(10.55), y, Inches(2.2), Inches(0.32),
                       f"{_fmt_brl(c['mensal'], 0)}  ·  {_fmt_pct(c['pct'])}",
                       size=13, color=MUTED, align=PP_ALIGN.RIGHT)
            y += Inches(0.62)
        self._rodape(s)

    def _slide_top_despesas(self):
        b = self.b
        top = (b.get('top_itens') or [])[:8]
        if not top:
            return
        s = self._slide()
        ano = self.r.get('ano_atual') or ''
        self._titulo(s, 'Maiores despesas do condomínio',
                     f'As 8 rubricas de maior peso na média mensal de {ano} '
                     f'· percentual = variação prevista da rubrica')

        nomes = [t['nome'][:38].upper() for t in reversed(top)]
        vals = [t['media_mensal'] / 1000.0 for t in reversed(top)]
        fig, ax = plt.subplots(figsize=(11.6, 4.6))
        barras = ax.barh(nomes, vals, color='#1A56DB', height=0.62, zorder=3)
        barras[-1].set_color('#0A1F44')  # a maior em destaque
        for sp in ax.spines.values():
            sp.set_visible(False)
        ax.tick_params(colors='#1B263B', labelsize=11.5, length=0)
        ax.xaxis.grid(True, color='#E4E9F4', linewidth=0.9)
        ax.set_axisbelow(True)
        ax.set_xlabel('R$ mil / mês', fontsize=11, color='#5A6B8C')
        def _var_rubrica(t):
            idx = t.get('indice')
            if idx is not None and 0.5 < idx < 2.0:
                return (idx - 1.0) * 100.0
            pn, md = t.get('prev_next_mensal'), t.get('media_mensal')
            if pn and md:
                return (pn / md - 1.0) * 100.0
            return None
        for barra, t in zip(barras, reversed(top)):
            var = _var_rubrica(t)
            etiqueta = _fmt_brl(t['media_mensal'], 0)
            if var is not None:
                etiqueta += f'  ·  {_fmt_pct(var, 1, sinal=True)}'
            ax.text(barra.get_width() + max(vals) * 0.012,
                    barra.get_y() + barra.get_height() / 2, etiqueta,
                    va='center', fontsize=9.5, color='#5A6B8C')
        fig.tight_layout()
        self._add_chart(s, fig, Inches(0.55), Inches(1.55), Inches(12.2))
        self._rodape(s)

    def _slide_inadimplencia(self):
        inad = self.extras.get('inadimplencia_total')
        acordos = self.extras.get('acordos_receber')
        if not inad and not acordos:
            return
        s = self._slide(bg=WHITE)
        self._titulo(s, 'Inadimplência e valores a receber',
                     'Pressão direta sobre o caixa do condomínio')

        cards = []
        if inad:
            cards.append(('Total em aberto', inad,
                          'Contas vencidas e a vencer em cobrança'))
        if acordos:
            cards.append(('Parcelas de acordo a receber', acordos,
                          'Renegociações em andamento'))
        w = Inches(5.9) if len(cards) == 2 else Inches(8.0)
        x0 = Inches(0.6) if len(cards) == 2 else Inches(2.65)
        for i, (rot, val, sub) in enumerate(cards):
            x = x0 + i * (w + Inches(0.33))
            self._card(s, x, Inches(2.0), w, Inches(2.3), fill=ICE)
            self._text(s, x + Inches(0.35), Inches(2.32), w - Inches(0.7),
                       Inches(0.35), rot.upper(), size=12.5, color=BLUE,
                       bold=True)
            self._text(s, x + Inches(0.35), Inches(2.75), w - Inches(0.7),
                       Inches(0.95), _fmt_brl(val, 2), size=40, color=RED,
                       bold=True)
            self._text(s, x + Inches(0.35), Inches(3.78), w - Inches(0.7),
                       Inches(0.4), sub, size=12, color=MUTED)

        taxa = self.r.get('taxa_ideal_mensal')
        if inad and taxa and taxa > 0:
            equivalente = inad / taxa
            self._rich(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(0.9), [
                ('O valor em aberto equivale a ', 16, TEXT, False),
                (f'{equivalente:,.0f} taxas condominiais'.replace(',', '.'),
                 16, RED, True),
                (' — despesas que precisam ser cobertas pelos condôminos adimplentes.',
                 16, TEXT, False),
            ])
        self._rodape(s)

    def _slide_previsao(self):
        b = self.b
        s = self._slide()
        ano = self.r.get('ano_proximo') or ''
        self._titulo(s, f'Previsão orçamentária {ano}',
                     'Composição do valor mensal a ratear entre as unidades')

        if b:
            linhas = [
                (f'Despesas previstas {ano}', b.get('prev_next_mensal')),
                (f'Fundo de Reserva ({_fmt_pct(self.r.get("fundo_reserva_pct"), 0)})',
                 b.get('fundo_mensal')),
            ]
            if b.get('garantidora_mensal'):
                linhas.append((
                    f'Garantidora ({_fmt_pct(b.get("garantidora_pct"), 1)})',
                    b['garantidora_mensal']))
            total = b.get('total_rateado_mensal')
            idx = b.get('indice_global')
        else:
            fr = (self.r.get('fundo_reserva') or 0) / 12.0 or None
            gar = (self.r.get('garantidora') or 0) / 12.0 or None
            linhas = [(f'Despesas previstas {ano}',
                       (self.r.get('total_despesas') or 0) / 12.0 or None),
                      (f'Fundo de Reserva ({_fmt_pct(self.r.get("fundo_reserva_pct"), 0)})', fr)]
            if gar:
                linhas.append((f'Garantidora ({_fmt_pct(self.r.get("garantidora_pct"), 0)})', gar))
            total = (self.r.get('total_rateado') or 0) / 12.0 or None
            idx = None

        y = Inches(1.9)
        for rot, val in linhas:
            self._card(s, Inches(0.6), y, Inches(7.6), Inches(0.78), fill=ICE_2)
            self._text(s, Inches(0.95), y + Inches(0.21), Inches(4.6),
                       Inches(0.36), rot, size=15, color=TEXT, bold=True)
            self._text(s, Inches(5.4), y + Inches(0.19), Inches(2.5),
                       Inches(0.4), _fmt_brl(val, 2), size=17, color=NAVY,
                       bold=True, align=PP_ALIGN.RIGHT)
            y += Inches(0.94)

        self._card(s, Inches(0.6), y + Inches(0.08), Inches(7.6), Inches(0.92),
                   fill=NAVY)
        self._text(s, Inches(0.95), y + Inches(0.35), Inches(4.6), Inches(0.4),
                   'TOTAL A RATEAR / MÊS', size=14, color=BLUE_LT, bold=True)
        self._text(s, Inches(5.4), y + Inches(0.3), Inches(2.5), Inches(0.46),
                   _fmt_brl(total, 2), size=20, color=WHITE, bold=True,
                   align=PP_ALIGN.RIGHT)

        # painel lateral: por unidade
        uni = self.r.get('num_unidades')
        fracoes = (b.get('fracoes') or []) if b else []
        por_fracao = not uni and b and b.get('taxa_media_provisionada')
        self._card(s, Inches(8.6), Inches(1.9), Inches(4.13),
                   Inches(5.0 if fracoes else 3.9), fill=ICE)
        self._text(s, Inches(8.95), Inches(2.25), Inches(3.4), Inches(0.4),
                   'RATEIO POR FRAÇÃO IDEAL' if por_fracao else 'RATEIO POR UNIDADE',
                   size=12.5, color=MUTED, bold=True)
        self._text(s, Inches(8.95), Inches(2.72), Inches(3.4), Inches(0.5),
                   'conforme blocos' if por_fracao else f'{uni or "—"} unidades',
                   size=17, color=TEXT, bold=True)
        self._text(s, Inches(8.95), Inches(3.5), Inches(3.4), Inches(0.4),
                   'Taxa média resultante' if por_fracao else 'Taxa condominial resultante',
                   size=12.5, color=MUTED)
        self._text(s, Inches(8.95), Inches(3.9), Inches(3.5), Inches(0.75),
                   _fmt_brl(self.r.get('taxa_ideal_mensal'), 2), size=31,
                   color=BLUE, bold=True)
        if fracoes:
            yf = Inches(4.85)
            for fr in fracoes[:4]:
                self._text(s, Inches(8.95), yf, Inches(2.5), Inches(0.3),
                           fr['nome'][:24], size=9.5, color=MUTED)
                self._text(s, Inches(11.15), yf, Inches(1.35), Inches(0.3),
                           _fmt_brl(fr['total'], 2), size=10, color=NAVY,
                           bold=True, align=PP_ALIGN.RIGHT)
                yf += Inches(0.42)
        elif idx:
            self._text(s, Inches(8.95), Inches(4.95), Inches(3.5), Inches(0.6),
                       f'Índice de reajuste aplicado: '
                       f'{_fmt_pct((idx - 1) * 100, 2, sinal=True)}',
                       size=11.5, color=MUTED)
        self._rodape(s)

    def _slide_proposta(self):
        taxa_nova = self.r.get('taxa_ideal_mensal')
        if not taxa_nova:
            return
        taxa_atual = self.r.get('taxa_atual')
        reajuste = self.r.get('reajuste_pct')
        s = self._slide()
        self._titulo(s, 'Proposta de reajuste da taxa condominial',
                     'Valor necessário para equilibrar as contas do próximo exercício')

        if taxa_atual:
            origem = self.b.get('taxa_atual_origem')
            if origem == 'aprovada':
                rotulo_atual = f'TAXA APROVADA {self.r.get("ano_atual") or ""}'.strip()
            elif origem == 'planilha':
                rotulo_atual = 'TAXA MÉDIA ATUAL'
            else:
                rotulo_atual = 'TAXA ATUAL'
            # taxa atual → taxa proposta
            self._card(s, Inches(0.9), Inches(2.2), Inches(4.7), Inches(2.6),
                       fill=ICE_2)
            self._text(s, Inches(1.25), Inches(2.6), Inches(4.0), Inches(0.4),
                       rotulo_atual, size=13, color=MUTED, bold=True)
            self._text(s, Inches(1.25), Inches(3.15), Inches(4.0), Inches(0.9),
                       _fmt_brl(taxa_atual, 2), size=38, color=MUTED, bold=True)
            self._text(s, Inches(1.25), Inches(4.2), Inches(4.0), Inches(0.4),
                       'por unidade / mês', size=12, color=MUTED)

            seta = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.85),
                                      Inches(3.25), Inches(1.35), Inches(0.55))
            seta.fill.solid()
            seta.fill.fore_color.rgb = BLUE
            seta.line.fill.background()
            seta.shadow.inherit = False

            self._card(s, Inches(7.45), Inches(2.0), Inches(5.0), Inches(3.0),
                       fill=NAVY)
            self._text(s, Inches(7.8), Inches(2.4), Inches(4.3), Inches(0.4),
                       'TAXA PROPOSTA', size=13, color=BLUE_LT, bold=True)
            self._text(s, Inches(7.8), Inches(2.95), Inches(4.4), Inches(1.0),
                       _fmt_brl(taxa_nova, 2), size=44, color=WHITE, bold=True)
            self._text(s, Inches(7.8), Inches(4.15), Inches(4.3), Inches(0.4),
                       'por unidade / mês', size=12, color=MUTED_ICE)

            if reajuste is not None:
                cor = RED if reajuste > 0 else GREEN
                pill = self._card(s, Inches(7.8), Inches(4.6), Inches(2.6),
                                  Inches(0.52), fill=WHITE, radius=0.5)
                self._text(s, Inches(7.8), Inches(4.7), Inches(2.6),
                           Inches(0.36), f'Reajuste de {_fmt_pct(abs(reajuste))}',
                           size=14, color=cor, bold=True,
                           align=PP_ALIGN.CENTER)
                dif = taxa_nova - taxa_atual
                self._text(s, Inches(0.9), Inches(5.5), Inches(11.5),
                           Inches(0.5),
                           f'Diferença de {_fmt_brl(abs(dif), 2)} por unidade/mês '
                           f'para cobrir o custo real de operação.',
                           size=14.5, color=TEXT)
        else:
            # sem taxa atual: hero único
            self._card(s, Inches(3.4), Inches(2.1), Inches(6.5), Inches(3.1),
                       fill=NAVY)
            self._text(s, Inches(3.4), Inches(2.55), Inches(6.5), Inches(0.4),
                       'TAXA CONDOMINIAL PROPOSTA', size=14, color=BLUE_LT,
                       bold=True, align=PP_ALIGN.CENTER)
            self._text(s, Inches(3.4), Inches(3.1), Inches(6.5), Inches(1.1),
                       _fmt_brl(taxa_nova, 2), size=52, color=WHITE, bold=True,
                       align=PP_ALIGN.CENTER)
            self._text(s, Inches(3.4), Inches(4.45), Inches(6.5), Inches(0.4),
                       'por unidade / mês', size=13, color=MUTED_ICE,
                       align=PP_ALIGN.CENTER)
        self._rodape(s)

    def _slide_conclusao(self):
        """
        Conclusão gerada a partir dos números calculados — texto template
        determinístico (sem geração livre), portanto sempre fiel aos dados.
        """
        b = self.b
        r = self.r
        pontos = []

        ano = r.get('ano_atual')
        ano_prox = r.get('ano_proximo')
        desvio = b.get('desvio_pct')
        if desvio is not None and b.get('realizado_medio_mensal'):
            if desvio > 1:
                pontos.append((
                    'Despesas acima do orçamento',
                    f'Em {ano}, a despesa média mensal foi de '
                    f'{_fmt_brl(b["realizado_medio_mensal"], 0)} — '
                    f'{_fmt_pct(desvio)} acima da previsão aprovada de '
                    f'{_fmt_brl(b.get("prev_aprovada_mensal"), 0)}.', RED))
            elif desvio < -1:
                pontos.append((
                    'Despesas dentro do orçamento',
                    f'Em {ano}, a despesa média mensal ficou '
                    f'{_fmt_pct(abs(desvio))} abaixo da previsão aprovada.',
                    GREEN))
            else:
                pontos.append((
                    'Execução alinhada ao orçamento',
                    f'Em {ano}, a despesa média mensal ficou em linha com a '
                    f'previsão aprovada (desvio de {_fmt_pct(desvio, 1, True)}).',
                    GREEN))

        inad = self.extras.get('inadimplencia_total')
        if inad:
            taxa = r.get('taxa_ideal_mensal')
            eq = ''
            if taxa:
                n = f'{inad / taxa:,.0f}'.replace(',', '.')
                eq = f', o equivalente a {n} taxas condominiais'
            pontos.append((
                'Inadimplência pressiona o caixa',
                f'Há {_fmt_brl(inad, 2)} em aberto{eq}. Esse custo é '
                f'absorvido pelos condôminos adimplentes.', AMBER))

        serie = b.get('serie_mensal') or []
        prev = b.get('prev_aprovada_mensal')
        if serie and prev:
            acima = sum(1 for p in serie if p['valor'] > prev)
            if acima >= len(serie) // 2:
                pontos.append((
                    'Tendência consistente, não pontual',
                    f'As despesas superaram a previsão em {acima} dos '
                    f'{len(serie)} meses do exercício — o desequilíbrio é '
                    f'estrutural, não eventual.', RED))

        taxa_nova = r.get('taxa_ideal_mensal')
        reajuste = r.get('reajuste_pct')
        if taxa_nova:
            det = ''
            if reajuste is not None:
                og = b.get('taxa_atual_origem')
                if og == 'aprovada':
                    origem = f'sobre a taxa aprovada em {ano}'
                elif og == 'planilha':
                    origem = 'sobre a taxa média atual'
                else:
                    origem = 'sobre a taxa vigente'
                det = (f' — reajuste de {_fmt_pct(abs(reajuste))} {origem}'
                       if reajuste >= 0 else
                       f' — redução de {_fmt_pct(abs(reajuste))} {origem}')
            pontos.append((
                'Recomendação',
                f'Aprovar a taxa condominial de {_fmt_brl(taxa_nova, 2)} por '
                f'unidade/mês para {ano_prox}{det}, garantindo o equilíbrio '
                f'financeiro do condomínio.', BLUE))

        if not pontos:
            return
        s = self._slide()
        self._titulo(s, 'Conclusão',
                     'Síntese dos indicadores apurados na análise')
        y = Inches(1.75)
        alt = Inches(1.14)
        for titulo, texto, cor in pontos[:4]:
            self._card(s, Inches(0.6), y, Inches(12.13), alt, fill=ICE_2)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95),
                                     y + Inches(0.24), Inches(0.16), Inches(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = cor
            dot.line.fill.background()
            dot.shadow.inherit = False
            self._text(s, Inches(1.32), y + Inches(0.14), Inches(11.0),
                       Inches(0.34), titulo, size=14.5, color=TEXT, bold=True)
            self._text(s, Inches(1.32), y + Inches(0.5), Inches(11.0),
                       Inches(0.56), texto, size=12.5, color=MUTED)
            y += alt + Inches(0.18)
        self._rodape(s)

    def _slide_encerramento(self):
        s = self._slide(bg=WHITE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.8), Inches(4.6),
                                  Inches(6.2), Inches(6.2))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ICE
        circ.line.fill.background()
        circ.shadow.inherit = False
        barra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W,
                                   Inches(0.22))
        barra.fill.solid()
        barra.fill.fore_color.rgb = NAVY
        barra.line.fill.background()
        barra.shadow.inherit = False

        taxa = self.r.get('taxa_ideal_mensal')
        ano = self.r.get('ano_proximo') or ''
        fracoes = (self.b.get('fracoes') or []) if self.b else []
        self._logo(s, Inches(0.6), Inches(0.55), altura=Inches(1.15), direita=True)
        self._text(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.5),
                   'PROPOSTA PARA APROVAÇÃO EM ASSEMBLEIA', size=15,
                   color=BLUE, bold=True)
        if taxa:
            rot_taxa = ('Taxa média condominial' if fracoes
                        else 'Taxa condominial')
            self._rich(s, Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.2), [
                (f'{rot_taxa} {ano}:  ', 32, NAVY, True),
                (_fmt_brl(taxa, 2), 42, BLUE, True),
                ('  / unidade / mês', 17, MUTED, False),
            ])
        if fracoes:
            y = Inches(3.85)
            self._text(s, Inches(0.9), y - Inches(0.42), Inches(11.5),
                       Inches(0.35), 'VALORES POR FRAÇÃO IDEAL', size=12.5,
                       color=MUTED, bold=True)
            for fr in fracoes[:5]:
                self._card(s, Inches(0.9), y, Inches(11.5), Inches(0.56),
                           fill=ICE_2, radius=0.22)
                self._text(s, Inches(1.2), y + Inches(0.12), Inches(8.6),
                           Inches(0.34), fr['nome'][:70], size=12, color=TEXT)
                self._text(s, Inches(9.9), y + Inches(0.1), Inches(2.2),
                           Inches(0.36), _fmt_brl(fr['total'], 2), size=14,
                           color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
                y += Inches(0.66)
        nome = self.r.get('nome_condominio')
        data = self.r.get('data_assembleia')
        rodape = ' · '.join(v for v in [nome, f'Assembleia em {data}' if data else None] if v)
        if rodape:
            self._text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
                       rodape, size=12.5, color=MUTED)

    def _slides_apendice(self):
        b = self.b
        itens = b.get('top_itens') or []
        if not itens:
            return
        ano = self.r.get('ano_atual') or ''
        por_pagina = 22
        paginas = [itens[i:i + por_pagina]
                   for i in range(0, len(itens), por_pagina)]
        for pg, bloco in enumerate(paginas, 1):
            s = self._slide()
            self._titulo(s, 'Apêndice — Detalhamento das rubricas',
                         f'Despesa média mensal de {ano} por rubrica '
                         f'(página {pg} de {len(paginas)})')
            col_w = Inches(5.95)
            metade = (len(bloco) + 1) // 2
            colunas = [bloco[:metade], bloco[metade:]]
            for c, coluna in enumerate(colunas):
                x = Inches(0.6) + c * (col_w + Inches(0.25))
                y = Inches(1.7)
                for i, it in enumerate(coluna):
                    if i % 2 == 0:
                        self._card(s, x, y - Inches(0.04), col_w,
                                   Inches(0.44), fill=ICE_2, radius=0.18)
                    self._text(s, x + Inches(0.18), y + Inches(0.05),
                               Inches(4.0), Inches(0.32), it['nome'][:42],
                               size=11, color=TEXT)
                    self._text(s, x + Inches(4.1), y + Inches(0.05),
                               Inches(1.68), Inches(0.32),
                               _fmt_brl(it['media_mensal'], 0), size=11,
                               color=MUTED, align=PP_ALIGN.RIGHT)
                    y += Inches(0.465)
            self._rodape(s)

    # ══════════════════════════════════════════════════════════════════
    def generate(self):
        # Ordem definida pelo cliente (Suport): detalhamento primeiro,
        # narrativa cresce até o panorama e fecha na proposta.
        self._slide_capa()
        self._slides_edital()
        self._slides_apendice()
        self._slide_top_despesas()
        self._slide_evolucao()
        self._slide_categorias()
        self._slide_panorama()
        self._slide_inadimplencia()
        self._slide_previsao()
        self._slide_encerramento()
        self._slide_proposta()
        self.prs.save(self.output_path)
        return self.output_path
